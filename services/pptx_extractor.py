from pptx import Presentation

def extract_pptx_slides(file_path):
    presentation = Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        slide_title = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            # Try to detect title
            if shape == slide.shapes.title:
                slide_title = text
            else:
                slide_text.append(text)

        slides_data.append({
            "source_type": "pptx",
            "page_number": None,
			"slide_number": idx,
			"title": slide_title,
			"text": "\n".join(slide_text)
		})

    return slides_data